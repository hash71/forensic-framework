"""Build the final 10-slide presentation deck from scratch.

Reuses the existing Project_Presentation.pptx as a template (preserves theme,
master, fonts, colours) but replaces every slide with content that matches the
final paper.
"""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

REPO = Path(__file__).resolve().parent.parent
SRC = REPO / "Project_Presentation.pptx"
DST = REPO / "paper_drafts" / "submission" / "Project_Presentation.pptx"
DST.parent.mkdir(parents=True, exist_ok=True)

# ----- palette -----
NAVY = RGBColor(0x1E, 0x3A, 0x8A)      # primary headers
SLATE = RGBColor(0x33, 0x33, 0x33)     # body text
MUTED = RGBColor(0x6B, 0x72, 0x80)     # captions
RED = RGBColor(0xC0, 0x39, 0x2B)       # failure markers
GREEN = RGBColor(0x16, 0xA3, 0x4A)     # success markers
BG = RGBColor(0xF5, 0xF5, 0xF5)        # subtle background


def _delete_all_slides(prs: Presentation) -> None:
    """Remove every existing slide while preserving masters / theme."""
    sldIdLst = prs.slides._sldIdLst
    rId_list = []
    for sldId in list(sldIdLst):
        rId_list.append(sldId.attrib[qn("r:id")])
        sldIdLst.remove(sldId)
    for rId in rId_list:
        prs.part.drop_rel(rId)


def _txt(tf, runs, *, size=18, bold=False, colour=SLATE, align=None, space_after=4):
    """Replace tf contents with one or more paragraph 'runs'.

    runs: list of (text, kwargs) tuples; each tuple becomes one paragraph.
          kwargs override the defaults for that paragraph only.
    """
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(runs):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align) or PP_ALIGN.LEFT
        p.space_after = Pt(opts.get("space_after", space_after))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        c = opts.get("colour", colour)
        if c is not None:
            run.font.color.rgb = c


def _add_textbox(slide, left_in, top_in, width_in, height_in, runs, **kw):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    _txt(box.text_frame, runs, **kw)
    box.text_frame.margin_left = Pt(0)
    box.text_frame.margin_right = Pt(0)
    box.text_frame.margin_top = Pt(2)
    box.text_frame.margin_bottom = Pt(2)
    return box


def _add_rule(slide, left_in, top_in, width_in, colour=NAVY, height_pt=2):
    """Thin horizontal rule for visual separation."""
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left_in), Inches(top_in),
                                 Inches(width_in), Pt(height_pt))
    box.fill.solid()
    box.fill.fore_color.rgb = colour
    box.line.fill.background()
    return box


def _new_slide(prs, layout_idx=0):
    """Add a blank slide; the source deck only ships one layout ('DEFAULT')."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # Remove any auto-populated placeholders from the layout so we can position
    # our own text boxes cleanly.
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


def _slide_header(slide, title, subtitle=None):
    """Standard slide header: title + optional subtitle + accent rule."""
    _add_textbox(slide, 0.5, 0.35, 12, 0.7,
                 [(title, dict(size=28, bold=True, colour=NAVY))])
    if subtitle:
        _add_textbox(slide, 0.5, 0.95, 12, 0.4,
                     [(subtitle, dict(size=14, colour=MUTED))])
    _add_rule(slide, 0.5, 1.35, 12.33, colour=NAVY, height_pt=1.5)


def slide_title(prs):
    s = _new_slide(prs)
    _add_textbox(s, 0.5, 2.4, 12.33, 1.3,
                 [("Augmenting Rule-Based Private-Cloud", dict(size=34, bold=True, colour=NAVY)),
                  ("Forensic Investigation with",            dict(size=34, bold=True, colour=NAVY)),
                  ("Evidence-Grounded LLM Reasoning",        dict(size=34, bold=True, colour=NAVY))])
    _add_rule(s, 0.5, 4.1, 12.33, colour=NAVY, height_pt=2)
    _add_textbox(s, 0.5, 4.4, 12.33, 0.6,
                 [("A Controlled Comparison of Rule-Based and LLM-Assisted Post-Incident Triage Under Citation Constraints",
                   dict(size=15, colour=MUTED))])
    _add_textbox(s, 0.5, 6.0, 12.33, 0.5,
                 [("Md. Nazmul Hasan", dict(size=20, bold=True, colour=SLATE))])
    _add_textbox(s, 0.5, 6.6, 12.33, 0.5,
                 [("May 2026", dict(size=14, colour=MUTED))])


def slide_problem_1(prs):
    s = _new_slide(prs)
    _slide_header(s, "The Problem (1/2): Private-Cloud Forensic Triage",
                  "No unified audit plane. Evidence is fragmented across many surfaces.")
    _add_textbox(s, 0.5, 1.7, 12, 0.5,
                 [("After an incident, the analyst must reconstruct the timeline from:", dict(size=16))])
    sources = [
        "Authentication logs (logins, failures, source IPs)",
        "File-access logs (reads, downloads, sizes)",
        "Administrative actions (privilege changes, log rotation)",
        "Network telemetry (DNS, firewall, VPN)",
        "Database query traces",
        "Web-server access logs",
        "Email metadata",
    ]
    body = [(f"•  {x}", dict(size=15, space_after=6)) for x in sources]
    _add_textbox(s, 0.7, 2.4, 12, 4.5, body)
    _add_textbox(s, 0.5, 6.5, 12, 0.6,
                 [("This work concerns post-incident triage, not real-time intrusion detection.",
                   dict(size=14, bold=True, colour=NAVY))])


def slide_problem_2(prs):
    s = _new_slide(prs)
    _slide_header(s, "The Problem (2/2): Rules vs LLMs",
                  "Two complementary failure modes; neither alone is sufficient.")

    # left column — rules
    _add_textbox(s, 0.5, 1.7, 6, 0.5,
                 [("Rule-Based Correlation", dict(size=18, bold=True, colour=NAVY))])
    rule_body = [
        ("✓ Deterministic, transparent, auditable", dict(size=14, space_after=6)),
        ("✓ Required in regulated environments", dict(size=14, space_after=6)),
        ("✗ Brittle threshold tuning", dict(size=14, colour=RED, space_after=6)),
        ("✗ High alert volume on legitimate-but-unusual activity", dict(size=14, colour=RED, space_after=6)),
        ("✗ Cannot aggregate across long time windows", dict(size=14, colour=RED, space_after=6)),
    ]
    _add_textbox(s, 0.7, 2.2, 6, 4, rule_body)

    # right column — LLM
    _add_textbox(s, 6.8, 1.7, 6, 0.5,
                 [("LLM-Assisted Reasoning", dict(size=18, bold=True, colour=NAVY))])
    llm_body = [
        ("✓ Synthesises weak signals across many events", dict(size=14, space_after=6)),
        ("✓ Reads role/ticket metadata as context", dict(size=14, space_after=6)),
        ("✓ Produces narrative, citable explanations", dict(size=14, space_after=6)),
        ("✗ Hallucinates events / actors / timelines", dict(size=14, colour=RED, space_after=6)),
        ("✗ Over-attributes intent and causality", dict(size=14, colour=RED, space_after=6)),
    ]
    _add_textbox(s, 7.0, 2.2, 6, 4, llm_body)

    _add_textbox(s, 0.5, 6.4, 12, 0.6,
                 [("Research question: can an evidence-grounded LLM, required to cite events and validated post-hoc, improve triage over rules alone?",
                   dict(size=14, bold=True, colour=NAVY))])


def slide_method(prs):
    s = _new_slide(prs)
    _slide_header(s, "Method: Nine-Stage Pipeline",
                  "Same normalised timeline feeds rule engine and LLM in parallel; validator inspects only the LLM output.")

    stages = [
        ("1. Scenario JSON", 0.5),
        ("2. Synthetic log generator", 1.45),
        ("3. Raw logs (7 sources)", 2.4),
        ("4. Source-specific parsers", 3.35),
        ("5. Normaliser → OCSF v1.1 / ECS 8.x", 4.3),
        ("6. Timeline reconstruction", 5.25),
        ("7. Correlation engine", 6.2),
    ]
    for label, top_in in stages:
        _add_textbox(s, 0.7, top_in, 5.8, 0.6,
                     [(label, dict(size=15, bold=True, colour=SLATE))])
        _add_rule(s, 0.7, top_in + 0.55, 5.8, colour=MUTED, height_pt=0.8)

    # Right side: parallel branches
    _add_textbox(s, 7.2, 0.5, 5.6, 0.5,
                 [("8. Parallel analysis", dict(size=15, bold=True, colour=NAVY))])
    _add_textbox(s, 7.2, 1.0, 5.6, 0.6,
                 [("Rule engine: 12 threshold-style rules (R001–R012)", dict(size=13, space_after=4)),
                  ("LLM analyst: Qwen 3.5‑27B; every claim must cite an event_id", dict(size=13))])

    _add_textbox(s, 7.2, 2.4, 5.6, 0.5,
                 [("9. Validation (LLM only)", dict(size=15, bold=True, colour=NAVY))])
    _add_textbox(s, 7.2, 2.9, 5.6, 1.4,
                 [("Seven checks: event-ID existence, chronology, actor consistency, entity consistency, volume claims, temporal claims, structural support",
                   dict(size=13))])

    _add_textbox(s, 7.2, 4.6, 5.6, 0.5,
                 [("10. Evaluation harness", dict(size=15, bold=True, colour=NAVY))])
    _add_textbox(s, 7.2, 5.1, 5.6, 1.5,
                 [("Verdict accuracy, F1 against attack steps, FP rate, validator counts — saved per scenario in evaluation_results.json",
                   dict(size=13))])

    _add_textbox(s, 0.5, 7.0, 12.33, 0.4,
                 [("Implementation: forensic-framework, commit 6c7ffbc; LLM endpoint-agnostic (Modal / vLLM / Ollama).",
                   dict(size=12, colour=MUTED))])


def slide_scenarios(prs):
    s = _new_slide(prs)
    _slide_header(s, "Four Primary Scenarios",
                  "Reproduce the original proposal scope; eleven extended scenarios add coverage and failure cases.")

    rows = [
        ("S1", "Normal baseline",  "BENIGN", "Control. Both systems should produce no alert."),
        ("S2", "Noisy benign (conference travel)", "BENIGN",
         "Senior dev abroad, late hours, elevated volume — false-positive resistance test."),
        ("S3", "Obvious external attack", "ATTACK",
         "Compromised credentials, off-hours foreign IP, privilege escalation, bulk download, log cleanup."),
        ("S4", "Slow multi-day insider", "ATTACK",
         "HR specialist incrementally widens access into finance and engineering across three days. "
         "No single day exceeds rule thresholds."),
    ]
    top = 1.65
    # Header row
    _add_textbox(s, 0.5, top, 0.7, 0.4, [("ID", dict(size=13, bold=True, colour=NAVY))])
    _add_textbox(s, 1.3, top, 4.0, 0.4, [("Scenario", dict(size=13, bold=True, colour=NAVY))])
    _add_textbox(s, 5.4, top, 1.4, 0.4, [("Truth", dict(size=13, bold=True, colour=NAVY))])
    _add_textbox(s, 6.9, top, 6.0, 0.4, [("Purpose", dict(size=13, bold=True, colour=NAVY))])
    _add_rule(s, 0.5, top + 0.4, 12.33, colour=NAVY, height_pt=1)

    cur = top + 0.55
    for sid, name, truth, purpose in rows:
        truth_colour = GREEN if truth == "BENIGN" else RED
        _add_textbox(s, 0.5, cur, 0.7, 1.1, [(sid, dict(size=14, bold=True, colour=SLATE))])
        _add_textbox(s, 1.3, cur, 4.0, 1.1, [(name, dict(size=13, colour=SLATE))])
        _add_textbox(s, 5.4, cur, 1.4, 1.1, [(truth, dict(size=13, bold=True, colour=truth_colour))])
        _add_textbox(s, 6.9, cur, 6.0, 1.1, [(purpose, dict(size=12, colour=SLATE))])
        _add_rule(s, 0.5, cur + 1.05, 12.33, colour=MUTED, height_pt=0.6)
        cur += 1.2

    _add_textbox(s, 0.5, 7.0, 12.33, 0.4,
                 [("Plus eleven extended scenarios (S5–S15) including two documented failure cases.",
                   dict(size=12, colour=MUTED))])


def slide_headline(prs):
    s = _new_slide(prs)
    _slide_header(s, "Headline Result",
                  "Same correlated timeline, two analysers; rule baseline is transparent, not optimal.")

    # Big numbers
    _add_textbox(s, 0.5, 1.7, 6, 0.6,
                 [("Rule engine", dict(size=18, bold=True, colour=SLATE))])
    _add_textbox(s, 0.5, 2.3, 6, 1.5,
                 [("10 / 15", dict(size=58, bold=True, colour=MUTED))])
    _add_textbox(s, 0.5, 4.0, 6, 0.6,
                 [("66.7% verdict accuracy", dict(size=16, colour=MUTED))])

    _add_textbox(s, 6.9, 1.7, 6, 0.6,
                 [("LLM (rule-context condition)", dict(size=18, bold=True, colour=NAVY))])
    _add_textbox(s, 6.9, 2.3, 6, 1.5,
                 [("14 / 15", dict(size=58, bold=True, colour=NAVY))])
    _add_textbox(s, 6.9, 4.0, 6, 0.6,
                 [("93.3% verdict accuracy", dict(size=16, bold=True, colour=NAVY))])

    _add_rule(s, 0.5, 4.9, 12.33, colour=NAVY, height_pt=1)

    bullets = [
        ("Validator: no invalid event-ID citations in 13 of 15 scenarios (S10, S13 cite a rule ID once each — 2/243 references = 0.8%).",
         dict(size=13, space_after=8)),
        ("Chronology violations: 2 of 15 (S4, S5: one out-of-order step each).",
         dict(size=13, space_after=8)),
        ("Entity-consistency violation: 1 of 15 (S12: IP named in narrative not present in cited events).",
         dict(size=13, space_after=8)),
        ("Rules misclassify on subtle multi-day attacks (S5, S10, S13) and on legitimate-after-hours maintenance (S6).",
         dict(size=13, space_after=8)),
    ]
    _add_textbox(s, 0.5, 5.0, 12.33, 2.5, bullets)


def slide_s14(prs):
    s = _new_slide(prs)
    _slide_header(s, "Failure Mode 1 — S14: Decoy Misdirection",
                  "Verdict-correct, suspect-wrong. Citation grounding does not catch wrong-suspect-with-real-evidence.")

    _add_textbox(s, 0.5, 1.7, 6, 0.5,
                 [("Setup", dict(size=16, bold=True, colour=NAVY))])
    setup = [
        ("user_04 (decoy): twenty-one rule alerts — brute-force, privilege escalation, bulk download.",
         dict(size=13, space_after=6)),
        ("user_02 (real attacker): quiet lateral movement and source-code exfiltration over DNS tunneling, internal IP.",
         dict(size=13, space_after=6)),
        ("Ground-truth attacker: user_02.",
         dict(size=13, bold=True, colour=NAVY)),
    ]
    _add_textbox(s, 0.7, 2.2, 6, 3, setup)

    _add_textbox(s, 6.9, 1.7, 6, 0.5,
                 [("LLM Output", dict(size=16, bold=True, colour=NAVY))])
    out = [
        ("Verdict: YES against ATTACK → binary verdict CORRECT.",
         dict(size=13, colour=GREEN, space_after=6)),
        ("Suspect: user_04 (the decoy) → attribution WRONG.",
         dict(size=13, colour=RED, space_after=6)),
        ("Every cited event is real and chronologically ordered.",
         dict(size=13, space_after=6)),
        ("Validator records zero violations.",
         dict(size=13, space_after=6)),
    ]
    _add_textbox(s, 7.0, 2.2, 6, 3, out)

    _add_rule(s, 0.5, 5.4, 12.33, colour=NAVY, height_pt=1)
    _add_textbox(s, 0.5, 5.5, 12.33, 1.8,
                 [("Lesson: citation grounding (every cited event exists) is one form of correctness; "
                   "interpretive grounding (cited evidence supports the chosen conclusion) is a different, stricter form. "
                   "An adversary placing a decoy in the evidence stream can satisfy citation grounding while still misdirecting the analyser.",
                   dict(size=13, space_after=4)),
                  ("Mitigation → multi-suspect hypothesis tracking: require the LLM to enumerate alternative actors and weigh evidence for each before committing to a primary suspect.",
                   dict(size=13, bold=True, colour=NAVY))])


def slide_s15(prs):
    s = _new_slide(prs)
    _slide_header(s, "Failure Mode 2 — S15: Rule-Context Amplification",
                  "Lone verdict failure. No-rule-context ablation reverses the verdict on all five runs.")

    _add_textbox(s, 0.5, 1.7, 6, 0.5,
                 [("Rule-context condition (headline)", dict(size=15, bold=True, colour=NAVY))])
    _add_textbox(s, 0.7, 2.2, 6, 2.5,
                 [("Input: timeline + 10 rule alerts (bulk_download, cross_department_access on quarter-close).",
                   dict(size=13, space_after=6)),
                  ("LLM verdict: YES on all 5 runs.", dict(size=14, bold=True, colour=RED, space_after=6)),
                  ("Ground truth: BENIGN (legitimate end-of-quarter activity).",
                   dict(size=13, space_after=6))])

    _add_textbox(s, 6.9, 1.7, 6, 0.5,
                 [("No-rule-context ablation", dict(size=15, bold=True, colour=NAVY))])
    _add_textbox(s, 7.0, 2.2, 6, 2.5,
                 [("Input: same timeline; rule-alert artefact removed from prompt.",
                   dict(size=13, space_after=6)),
                  ("LLM verdict: NO on all 5 runs.", dict(size=14, bold=True, colour=GREEN, space_after=6)),
                  ("Same model, same temperature (0.1), same prompt template otherwise.",
                   dict(size=13, space_after=6))])

    _add_rule(s, 0.5, 4.9, 12.33, colour=NAVY, height_pt=1)
    _add_textbox(s, 0.5, 5.0, 12.33, 2.3,
                 [("Verdict reverses unanimously on the presence of the rule-alert artefact alone.",
                   dict(size=14, bold=True, colour=NAVY, space_after=8)),
                  ("Reading: the LLM treats rule alerts as quasi-ground-truth signals, not as candidate evidence to weigh — ten alerts push the posterior past the verdict threshold.",
                   dict(size=13, space_after=8)),
                  ("Mitigation → calibrated rule-context (confidence-weighted alerts), two-stage prompting (score alerts, then verdict), counterfactual ablation as a deployed check on YES verdicts.",
                   dict(size=13))])


def slide_limits(prs):
    s = _new_slide(prs)
    _slide_header(s, "Limitations and Threats to Validity",
                  "What these results do not establish.")

    items = [
        ("Internal", "Author-designed scenarios, author-tuned prompt; mitigated by publishing both and keeping S14/S15 in the corpus.",
         dict(size=14, bold=True, colour=NAVY)),
        ("External", "Synthetic data; one modelled organisation (Dhaka, four user roles). Real enterprise logs are noisier, sparser in ground-truth labelling.",
         dict(size=14, bold=True, colour=NAVY)),
        ("Construct", "Verdict accuracy is a coarse construct. Attack-step F1 disagrees with verdict on S8 (right answer, wrong vocabulary). No human-analyst usefulness study.",
         dict(size=14, bold=True, colour=NAVY)),
        ("Validator scope", "Catches event-ID, chronology, actor, entity, structural-support violations. Does NOT catch causal overclaim, missing alternatives, or wrong-suspect-with-grounded-evidence (S14).",
         dict(size=14, bold=True, colour=NAVY)),
        ("Reproducibility", "Single representative run for the 14 non-S15 scenarios; only S15 has N=5. Multi-model and run-variance characterisation are future work.",
         dict(size=14, bold=True, colour=NAVY)),
    ]
    cur = 1.7
    for label, body, label_opts in items:
        _add_textbox(s, 0.5, cur, 2.5, 0.5, [(label, label_opts)])
        _add_textbox(s, 3.1, cur, 9.7, 1.0, [(body, dict(size=13, space_after=4))])
        cur += 0.95


def slide_conclusion(prs):
    s = _new_slide(prs)
    _slide_header(s, "Conclusion and Future Directions",
                  "Bounded contribution; clear next steps.")

    _add_textbox(s, 0.5, 1.7, 12.33, 1.7,
                 [("Evidence-grounded LLM reasoning improves post-incident forensic triage over a transparent rule baseline (10/15 → 14/15) "
                   "while producing no invalid event-ID citations in 13 of 15 scenarios. Two failure modes bound the contribution: "
                   "S14 (decoy misdirection — verdict-correct, suspect-wrong) and S15 (rule-context amplification, ablation-confirmed).",
                   dict(size=13, space_after=8)),
                  ("Appropriate role: analyst-supervised triage assistant. Not autonomous detection, not a final attribution authority.",
                   dict(size=13, bold=True, colour=NAVY))])

    _add_rule(s, 0.5, 3.6, 12.33, colour=NAVY, height_pt=1)
    _add_textbox(s, 0.5, 3.7, 12.33, 0.5,
                 [("Future Directions", dict(size=16, bold=True, colour=NAVY))])

    futures = [
        "Multi-model evaluation (LLaMA, Mistral, Phi-4) under identical prompts and validator.",
        "Statistical rigour: N=5 runs per scenario, mean ± std on verdict accuracy and F1.",
        "Adversarial robustness extending S14: vary decoy:real ratio, decoy actor count, decoy event volume.",
        "Human-analyst usefulness study (correctness, evidence clarity, overclaim risk).",
        "Real-world dataset run (LANL, DARPA OpTC, anonymised institutional logs).",
        "Stronger validator: causal-overclaim detection, alternative-explanation enumeration, multi-suspect tracking.",
    ]
    body = [(f"•  {x}", dict(size=13, space_after=4)) for x in futures]
    _add_textbox(s, 0.7, 4.3, 12, 3, body)


def main() -> None:
    prs = Presentation(str(SRC))
    _delete_all_slides(prs)

    slide_title(prs)
    slide_problem_1(prs)
    slide_problem_2(prs)
    slide_method(prs)
    slide_scenarios(prs)
    slide_headline(prs)
    slide_s14(prs)
    slide_s15(prs)
    slide_limits(prs)
    slide_conclusion(prs)

    prs.save(str(DST))
    print(f"saved -> {DST}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
